import sys
import subprocess
from pathlib import Path
from openai import OpenAI
from dotenv import load_dotenv
import os
import argparse
import logging
import json
import urllib.parse
import requests
from typing import Dict, List, Optional

# ============ CONFIGURATION ============
# Load environment variables first
load_dotenv()

# Get API key and validate
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError(
        "OPENAI_API_KEY environment variable is not set. Please create a .env file with your API key."
    )

# Initialize OpenAI client after validating API key
client = OpenAI(api_key=api_key)


# Load system prompt
def load_system_prompt() -> str:
    """Load the system prompt from system_prompt.txt."""
    try:
        with open("system_prompt.txt", "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        logging.error(f"Error loading system prompt: {e}")
        return "You are an assistant that analyzes lecture content and extracts a structured topic hierarchy."


# ============ LOGGING SETUP ============
def setup_logging(verbosity: int, log_file: str = None):
    """Set up logging based on verbosity level and optional log file."""
    # Remove any existing handlers
    for handler in logging.root.handlers[:]:
        logging.root.removeHandler(handler)

    # Set up logging format
    formatter = logging.Formatter("%(asctime)s - %(levelname)s - %(message)s")

    # Set root logger level based on verbosity
    if verbosity == 0:
        logging.root.setLevel(logging.ERROR)
    elif verbosity == 1:
        logging.root.setLevel(logging.INFO)
    elif verbosity == 2:
        logging.root.setLevel(logging.DEBUG)

    # Set up console handler with the same level as root logger
    console_handler = logging.StreamHandler()
    console_handler.setLevel(logging.root.level)
    console_handler.setFormatter(formatter)
    logging.root.addHandler(console_handler)

    # Set up file handler if log file is specified
    if log_file:
        file_handler = logging.FileHandler(log_file)
        file_handler.setLevel(logging.root.level)  # Use same level as root logger
        file_handler.setFormatter(formatter)
        logging.root.addHandler(file_handler)


# ============ FILE HANDLERS ============


def check_dependencies():
    """Check if required dependencies are installed."""
    try:
        import pypdf
        import pptx
    except ImportError as e:
        logging.error(f"Required dependency not installed: {e}")
        logging.info("Please install required packages: pip install pypdf python-pptx")
        sys.exit(1)


def extract_text_from_pdf(pdf_path: Path) -> str:
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(pdf_path))
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    except Exception as e:
        logging.error(f"Error processing PDF {pdf_path.name}: {e}")
        return ""


def extract_text_from_pptx(pptx_path: Path) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        logging.error(f"Error processing PPTX {pptx_path.name}: {e}")
        return ""


def transcribe_audio_with_whisper(mp4_path: Path, verbosity: int = 1) -> str:
    try:
        logging.info(f"Transcribing {mp4_path.name} using Whisper...")

        # Capture output if verbosity is not 2
        capture_output = verbosity < 2

        result = subprocess.run(
            [
                "whisper",
                str(mp4_path),
                "--model",
                "base",
                "--language",
                "en",
                "--output_format",
                "txt",
            ],
            check=True,
            capture_output=capture_output,
            text=True,
        )

        # If we captured output and verbosity is 2, print it
        if capture_output and verbosity >= 2:
            print(result.stdout)

        transcript_file = mp4_path.with_suffix(".txt")
        return transcript_file.read_text() if transcript_file.exists() else ""
    except subprocess.CalledProcessError as e:
        logging.error(f"Error running Whisper on {mp4_path.name}: {e}")
        return ""
    except Exception as e:
        logging.error(f"Error processing MP4 {mp4_path.name}: {e}")
        return ""


# ============ TEXT COMBINATION ============


class OpenAIHelper:
    """Handles OpenAI API interactions and token calculations."""

    def __init__(self):
        # Current OpenAI pricing (as of 2024)
        self.PRICING = {
            "gpt-4": {
                "input": 0.03 / 1000,  # $0.03 per 1K tokens
                "output": 0.06 / 1000,  # $0.06 per 1K tokens
            },
            "gpt-3.5-turbo": {
                "input": 0.001 / 1000,  # $0.001 per 1K tokens
                "output": 0.002 / 1000,  # $0.002 per 1K tokens
            },
        }
        self.system_prompt = load_system_prompt()

    def calculate_token_cost(
        self, model: str, prompt_tokens: int, completion_tokens: int
    ) -> float:
        """Calculate the cost of an API call based on token usage."""
        if model not in self.PRICING:
            return 0.0

        input_cost = prompt_tokens * self.PRICING[model]["input"]
        output_cost = completion_tokens * self.PRICING[model]["output"]
        return input_cost + output_cost


class TextProcessor:
    """Handles text processing operations."""

    def __init__(self, verbosity: int = 1):
        self.verbosity = verbosity
        self.logger = LoggingHelper(verbosity)
        self.openai = OpenAIHelper()


class DocumentProcessor:
    """Handles processing of individual documents and their content."""

    def __init__(self, verbosity: int = 1):
        self.verbosity = verbosity
        self.logger = LoggingHelper(verbosity)
        self.openai = OpenAIHelper()
        self.existing_topics = set()

    def process_directory(self, directory: Path) -> Dict:
        """Process all documents in a directory and return combined topics."""
        all_topics = {"topics": {}, "associated_topics": {}}
        total_cost = 0.0

        for file in directory.iterdir():
            if file.suffix.lower() in [".pdf", ".pptx", ".ppt", ".mp4"]:
                try:
                    self.logger.log_subsection(f"Processing: {file.name}")
                    topics, cost = self.process_single_document(file)
                    if topics:
                        # Update existing topics set
                        if "topics" in topics:
                            self.existing_topics.update(topics["topics"].keys())
                        self._merge_topics(all_topics, topics)
                        total_cost += cost
                except Exception as e:
                    logging.error(f"Error processing {file.name}: {e}")
                    continue

        if self.verbosity >= 1:
            logging.info(f"Total cost for all documents: ${total_cost:.4f}\n")

        return all_topics

    def process_single_document(self, file: Path) -> tuple[Dict, float]:
        """Process a single document and extract its topics."""
        ext = file.suffix.lower()
        content = ""

        # Extract content based on file type
        if ext == ".pdf":
            self.logger.log_subsection(f"Processing PDF: {file.name}")
            content = extract_text_from_pdf(file)
        elif ext in [".pptx", ".ppt"]:
            self.logger.log_subsection(f"Processing PowerPoint: {file.name}")
            content = extract_text_from_pptx(file)
        elif ext == ".mp4":
            self.logger.log_subsection(f"Processing MP4: {file.name}")
            content = transcribe_audio_with_whisper(file, self.verbosity)

        if not content:
            return {}, 0.0

        try:
            # Prepare existing topics information
            existing_topics_info = ""
            if self.existing_topics:
                existing_topics_info = "\n\nExisting Topics:\n" + "\n".join(
                    f"- {topic}" for topic in sorted(self.existing_topics)
                )

            # Process content with GPT-4
            response = client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {
                        "role": "system",
                        "content": self.openai.system_prompt,
                    },
                    {
                        "role": "user",
                        "content": f"Extract topics and their structure from the following content:{existing_topics_info}\n\n{content}",
                    },
                ],
                temperature=0.4,
            )

            # Calculate cost
            prompt_tokens = response.usage.prompt_tokens
            completion_tokens = response.usage.completion_tokens
            cost = self.openai.calculate_token_cost(
                "gpt-4", prompt_tokens, completion_tokens
            )

            if self.verbosity >= 1:
                self.logger.log_token_usage(
                    prompt_tokens, completion_tokens, cost, "GPT-4"
                )

            # Parse the JSON response
            try:
                topics = json.loads(response.choices[0].message.content)
                return topics, cost
            except json.JSONDecodeError as e:
                logging.error(f"Error parsing topics JSON: {e}")
                return {}, cost

        except Exception as e:
            logging.error(f"Error processing document: {e}")
            return {}, 0.0

    def _merge_topics(self, all_topics: Dict, new_topics: Dict) -> None:
        """Merge topics from a document into the main topics dictionary."""
        if not new_topics.get("topics"):
            return

        # Ensure all_topics has the required structure
        if "topics" not in all_topics:
            all_topics["topics"] = {}
        if "associated_topics" not in all_topics:
            all_topics["associated_topics"] = {}

        # Merge topics and subtopics
        for topic, subtopics in new_topics["topics"].items():
            if topic not in all_topics["topics"]:
                all_topics["topics"][topic] = subtopics
            else:
                # Merge subtopics
                self._merge_subtopics(all_topics["topics"][topic], subtopics)

        # Merge associated topics
        if "associated_topics" in new_topics:
            for topic, associated in new_topics["associated_topics"].items():
                if topic not in all_topics["associated_topics"]:
                    all_topics["associated_topics"][topic] = []
                all_topics["associated_topics"][topic].extend(associated)
                # Remove duplicates while preserving order
                all_topics["associated_topics"][topic] = list(
                    dict.fromkeys(all_topics["associated_topics"][topic])
                )

    def _merge_subtopics(self, existing: Dict, new: Dict) -> None:
        """Recursively merge subtopics."""
        for subtopic, content in new.items():
            if subtopic not in existing:
                existing[subtopic] = content
            elif isinstance(content, dict) and isinstance(existing[subtopic], dict):
                self._merge_subtopics(existing[subtopic], content)


class LoggingHelper:
    """Handles all logging operations with consistent formatting."""

    def __init__(self, verbosity: int = 1):
        self.verbosity = verbosity

    def log_section_header(self, title: str, width: int = 80) -> None:
        """Log a section header with visual separators."""
        separator = "=" * width
        logging.info(f"\n{separator}")
        logging.info(title)
        logging.info(f"{separator}\n")

    def log_subsection(self, title: str, width: int = 40) -> None:
        """Log a subsection header with visual separators."""
        logging.info(f"\n{title}:")
        logging.info("-" * width)

    def log_token_usage(
        self, prompt_tokens: int, completion_tokens: int, cost: float, model: str = None
    ) -> None:
        """Log token usage and cost information."""
        self.log_subsection("Token Usage")
        if model:
            logging.info(f"Model: {model}")
        logging.info(f"Prompt tokens: {prompt_tokens}")
        logging.info(f"Completion tokens: {completion_tokens}")
        logging.info(f"Total cost: ${cost:.4f}\n")

    def log_processing_summary(
        self, duplicates_found: int, original_count: int, unique_count: int
    ) -> None:
        """Log a summary of the processing results."""
        self.log_subsection("Processing Summary")
        logging.info(f"Duplicates merged: {duplicates_found}")
        logging.info(f"Original topics: {original_count}")
        logging.info(f"Unique topics: {unique_count}\n")

        if self.verbosity >= 2:
            logging.debug("\nFinal Topic Structure:")
            logging.debug("-" * 40)

    def log_duplicate_found(
        self, title: str, original_length: int, new_length: int, combined_length: int
    ) -> None:
        """Log information about a found duplicate."""
        self.log_subsection("Duplicate Found")
        logging.info(f"Topic: {title}")
        logging.info(f"Original content length: {original_length} chars")
        logging.info(f"New content length: {new_length} chars")
        logging.info(f"Combined content length: {combined_length} chars\n")


class TopicProcessor:
    """Handles topic processing and structure generation."""

    def __init__(self, verbosity: int = 1):
        self.verbosity = verbosity
        self.logger = LoggingHelper(verbosity)

    def process_topics(self, topics: str) -> Dict:
        """Process topics into a structured JSON format."""
        self.logger.log_section_header("Processing Topics")

        try:
            # Parse the topics string into a structured format
            topic_structure = {}
            associated_topics = {}
            current_topic = None
            current_subtopics = {}
            current_level = 0
            current_path = []

            lines = topics.split("\n")
            for line in lines:
                line = line.strip()
                if not line or line == "---":
                    continue

                # Count leading # to determine level
                level = 0
                while line.startswith("#"):
                    level += 1
                    line = line[1:].strip()

                if level == 0:
                    # This is a main topic
                    if current_topic:
                        topic_structure[current_topic] = current_subtopics
                    current_topic = line
                    current_subtopics = {}
                    current_path = [current_topic]
                    associated_topics[current_topic] = []
                else:
                    # This is a subtopic
                    if level <= current_level:
                        # Go back up the tree
                        current_path = current_path[: level - 1]

                    current_path.append(line)
                    current_level = level

                    # Add to the structure
                    current = topic_structure
                    for i, path in enumerate(current_path[:-1]):
                        if path not in current:
                            current[path] = {}
                        current = current[path]
                    current[current_path[-1]] = {}

            # Add the last topic
            if current_topic:
                topic_structure[current_topic] = current_subtopics

            # Get associated topics from the LLM
            for topic in topic_structure.keys():
                associated = self._get_associated_topics(
                    topic, list(topic_structure.keys())
                )
                associated_topics[topic] = associated

            return {"topics": topic_structure, "associated_topics": associated_topics}

        except Exception as e:
            self.logger.log_subsection("Error Processing Topics")
            logging.error(f"Error processing topics: {e}")
            return {"topics": {}, "associated_topics": {}}

    def _get_associated_topics(self, topic: str, all_topics: List[str]) -> List[str]:
        """Get associated topics for a given topic using GPT-3.5-turbo."""
        try:
            prompt = f"""Given the topic "{topic}" and the following list of topics:
{", ".join(all_topics)}

Return a JSON array of the 3 most closely related topics from the list. Only include topics that are actually in the list."""

            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {
                        "role": "system",
                        "content": "You are a helpful assistant that identifies related topics.",
                    },
                    {"role": "user", "content": prompt},
                ],
                temperature=0.3,
            )

            # Parse the response and ensure it's a valid list of topics
            try:
                associated = json.loads(response.choices[0].message.content)
                # Filter to ensure only valid topics are included
                return [t for t in associated if t in all_topics][:3]
            except:
                return []

        except Exception as e:
            logging.error(f"Error getting associated topics for {topic}: {e}")
            return []


# ============ NOTE CREATION ============


def find_obsidian_vault_root(path: Path) -> Path:
    """Find the Obsidian vault root directory by looking for .obsidian/ directory."""
    current = path
    while current != current.parent:  # Stop at root directory
        if (current / ".obsidian").exists():
            return current
        current = current.parent
    return None


def generate_tag_path(vault_root: Path, current_path: Path) -> str:
    """Generate the tag path from vault root to current directory."""
    if not vault_root:
        return ""
    # Get relative path from vault root
    rel_path = current_path.relative_to(vault_root)
    # Convert to string and remove whitespace
    tag_path = str(rel_path).replace(" ", "")
    return tag_path


def create_obsidian_notes(
    topics: str, output_dir: Path, interactive: bool = False
) -> None:
    """Create Obsidian notes from topics dictionary."""
    # Ensure output directory exists
    output_dir.mkdir(parents=True, exist_ok=True)
    logging.debug(f"Output directory: {output_dir}")

    # Find Obsidian vault root and generate tag path
    vault_root = find_obsidian_vault_root(output_dir)
    tag_path = generate_tag_path(vault_root, output_dir)
    base_tags = ["todo"]
    if tag_path:
        base_tags.append(tag_path)
    tags_str = ", ".join(base_tags)
    logging.debug(f"Tags: {tags_str}")

    # Parse topics from JSON if it's a string
    if isinstance(topics, str):
        try:
            topics_dict = json.loads(topics)
        except json.JSONDecodeError as e:
            logging.error(f"Error parsing topics JSON: {e}")
            return
    else:
        topics_dict = topics

    # Create notes for each topic
    for topic, subtopics in topics_dict["topics"].items():
        try:
            # Format filename
            filename = "".join(c for c in topic if c.isalnum() or c in "._- ")
            if not filename.endswith(".md"):
                filename += ".md"

            if interactive:
                print(f"\nFile to be created: {filename}")
                while True:
                    response = input("Create this file? (y/n/a): ").lower()
                    if response == "y":
                        break
                    elif response == "n":
                        logging.info(f"Skipping file: {filename}")
                        continue
                    elif response == "a":
                        logging.info("Aborting all file creation")
                        return
                    else:
                        print("Invalid input. Please enter y, n, or a.")

            # Create the file
            filepath = output_dir / filename
            with open(filepath, "w", encoding="utf-8") as f:
                # Write YAML frontmatter
                f.write("---\n")
                f.write(f"tags: {tags_str}\n")
                f.write("---\n\n")

                # Format and write content
                note_formatter = NoteFormatter(verbosity=1)
                content = note_formatter.format_topic_content(
                    topic, subtopics, topics_dict["associated_topics"].get(topic, [])
                )
                f.write(content)

            logging.info(f"Created: {filepath.name}")
        except Exception as e:
            logging.error(f"Error creating file {filename}: {e}")
            logging.error(f"Error details: {str(e)}")
            continue


# ============ MAIN ============


def main():
    """Main function to process documents and create notes."""
    parser = argparse.ArgumentParser(
        description="Generate structured notes from various learning resources."
    )
    parser.add_argument("resource_path", help="Path to the resource folder")
    parser.add_argument(
        "-v0", action="store_const", const=0, dest="verbosity", help="Silent mode"
    )
    parser.add_argument(
        "-v1",
        action="store_const",
        const=1,
        dest="verbosity",
        help="Basic progress output",
    )
    parser.add_argument(
        "-v2", action="store_const", const=2, dest="verbosity", help="Detailed output"
    )
    parser.add_argument("-l", "--log-file", help="Log to file")
    parser.add_argument(
        "-i", "--interactive", action="store_true", help="Interactive mode"
    )
    parser.set_defaults(verbosity=1)

    args = parser.parse_args()
    setup_logging(args.verbosity, args.log_file)

    # Check dependencies first
    check_dependencies()

    # Convert to Path object
    resource_path = Path(args.resource_path)
    root_path = resource_path.parent

    if not resource_path.exists():
        logging.error(f"{resource_path} does not exist.")
        sys.exit(1)

    # Initialize processors
    doc_processor = DocumentProcessor(args.verbosity)
    note_formatter = NoteFormatter(args.verbosity)
    logger = LoggingHelper(args.verbosity)

    logger.log_section_header("Processing Documents")

    # Process all documents and get combined topics
    topics = doc_processor.process_directory(resource_path)
    if not topics:
        logging.error("No topics were generated.")
        sys.exit(1)

    # Create notes directory
    notes_dir = root_path / "notes"
    notes_dir.mkdir(exist_ok=True)

    # Create notes for each topic
    logger.log_section_header("Creating Notes")
    for topic, subtopics in topics["topics"].items():
        try:
            # Format the note content
            content = note_formatter.format_topic_content(
                topic, subtopics, topics["associated_topics"].get(topic, [])
            )

            # Create the note file
            note_file = notes_dir / f"{topic}.md"
            with open(note_file, "w", encoding="utf-8") as f:
                f.write(content)

            if args.verbosity >= 1:
                logging.info(f"Created note: {note_file.name}")

        except Exception as e:
            logging.error(f"Error creating note for {topic}: {e}")
            continue

    logging.info("Done!")


if __name__ == "__main__":
    main()


class NoteFormatter:
    """Handles formatting of notes based on topic structure."""

    def __init__(self, verbosity: int = 1):
        self.verbosity = verbosity
        self.logger = LoggingHelper(verbosity)

    def get_wikipedia_definition(self, topic: str) -> str:
        """Get the definition of a topic from Wikipedia."""
        try:
            # First search for the topic to find the closest match
            search_url = f"https://en.wikipedia.org/w/api.php?action=query&list=search&srsearch={urllib.parse.quote(topic)}&format=json"
            search_response = requests.get(search_url)
            search_response.raise_for_status()
            search_data = search_response.json()

            # If no search results found, return blank definition
            if not search_data.get("query", {}).get("search"):
                return ""

            # Get the title of the closest match
            closest_match = search_data["query"]["search"][0]["title"]

            # Now get the summary for the closest match
            summary_url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{urllib.parse.quote(closest_match)}"
            summary_response = requests.get(summary_url)
            summary_response.raise_for_status()
            summary_data = summary_response.json()

            return summary_data.get("extract", "")

        except Exception as e:
            self.logger.log_subsection("Wikipedia API Error")
            logging.error(f"Error fetching definition for {topic}: {e}")
            return ""

    def format_topic_content(
        self, topic: str, subtopics: Dict, associated_topics: List[str]
    ) -> str:
        """Format the content for a topic file."""
        # Get Wikipedia definition
        definition = self.get_wikipedia_definition(topic)

        # Start building the content
        content = [f"# {topic}", "", "## Definition", f"*{definition}*", ""]

        # Add subtopics with proper heading levels
        self._add_subtopics(content, subtopics, 2)

        # Add Related Topics section
        if associated_topics:
            content.extend(
                [
                    "",
                    "## Related Topics",
                    *[
                        f"- [{topic}]({urllib.parse.quote(topic)}.md)"
                        for topic in associated_topics
                    ],
                ]
            )

        return "\n".join(content)

    def _add_subtopics(self, content: List[str], subtopics: Dict, level: int) -> None:
        """Recursively add subtopics with proper heading levels."""
        for subtopic, nested in subtopics.items():
            # Add the current subtopic
            content.append(f"{'#' * level} {subtopic}")
            content.append("")

            # If there are nested subtopics, add them
            if isinstance(nested, dict):
                self._add_subtopics(content, nested, level + 1)
            elif isinstance(nested, str):
                # If it's a string, it's the content
                content.append(nested)
                content.append("")
