import pytest
import os
import sys
from pathlib import Path
import tempfile
import shutil
import logging
from unittest.mock import patch, MagicMock, call
import json
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
from src.document_processor import DocumentProcessor
from src.openai_helper import OpenAIHelper
from src.utils import (
    extract_text_from_pdf,
    extract_text_from_pptx,
    transcribe_audio_with_whisper,
    extract_text_from_txt,
)
from src.logging_helper import setup_colored_logging

# Add the project root to the Python path
sys.path.insert(0, str(Path(__file__).parent.parent))

from src import (
    NoteFormatter,
    LoggingHelper,
    setup_logging,
    check_dependencies,
    main,
)


@pytest.fixture
def temp_dir():
    """Create a temporary directory for testing."""
    with tempfile.TemporaryDirectory() as tmpdirname:
        yield Path(tmpdirname)


@pytest.fixture
def sample_pdf(tmp_path):
    """Create a sample PDF file for testing."""
    pdf_path = tmp_path / "sample.pdf"
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Sample PDF content", ln=1, align="C")
    pdf.output(str(pdf_path))
    return pdf_path


@pytest.fixture
def sample_pptx(tmp_path):
    """Create a sample PPTX file for testing."""
    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Sample PPTX content"
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def sample_pdf_text():
    """Sample text content from PDF."""
    return "Sample PDF content"


@pytest.fixture
def learning_resources_dir(temp_dir):
    """Create a learning resources directory for testing."""
    resources_dir = temp_dir / "Learning Resources"
    resources_dir.mkdir()
    return resources_dir


@pytest.fixture(autouse=True)
def mock_openai_env(monkeypatch):
    """Mock OpenAI environment variables."""
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")


@pytest.fixture
def mock_openai_client():
    """Create a mock OpenAI client."""
    with patch("openai.OpenAI") as mock_client:
        mock_instance = MagicMock()
        mock_instance.chat.completions.create.return_value = MagicMock(
            choices=[
                MagicMock(
                    message=MagicMock(
                        content='{"topics": {"Test Topic": {"Subtopic": "Description"}}, "associated_topics": {"Test Topic": ["Related"]}}'
                    )
                )
            ],
            usage=MagicMock(prompt_tokens=100, completion_tokens=50),
        )
        mock_client.return_value = mock_instance
        yield mock_instance


@pytest.fixture
def mock_openai_response():
    """Create a mock OpenAI response."""
    mock_response = MagicMock()
    mock_response.choices = [
        MagicMock(
            message=MagicMock(
                content='{"topics": {"Test Topic": {"Subtopic": "Description"}}, "associated_topics": {"Test Topic": ["Related"]}}'
            )
        )
    ]
    mock_response.usage = MagicMock(prompt_tokens=100, completion_tokens=50)
    return mock_response


@pytest.fixture
def mock_openai_helper(mock_openai_client):
    """Create a mock OpenAIHelper."""
    with patch("openai.OpenAI") as mock_openai_class:
        mock_openai_class.return_value = mock_openai_client
        helper = OpenAIHelper()
        helper.client = mock_openai_client
        return helper


@pytest.fixture
def mock_logger():
    """Create a mock LoggingHelper."""
    return LoggingHelper(verbosity=1)


@pytest.fixture
def document_processor(mock_openai_helper, mock_logger):
    """Create a DocumentProcessor with mocked dependencies."""
    return DocumentProcessor(
        verbosity=1, openai_helper=mock_openai_helper, logger=mock_logger
    )


@pytest.fixture
def note_formatter():
    """Create a NoteFormatter instance."""
    return NoteFormatter(verbosity=1)


@pytest.fixture
def sample_mp4(tmp_path):
    """Create a sample MP4 file for testing."""
    mp4_path = tmp_path / "test.mp4"
    mp4_path.write_text("Test video content")
    return mp4_path


@pytest.fixture
def sample_txt(tmp_path):
    """Create a sample text file for testing."""
    txt_path = tmp_path / "sample.txt"
    txt_path.write_text("Sample text content")
    return txt_path


def test_extract_text_from_pdf(sample_pdf, sample_pdf_text):
    """Test PDF text extraction."""
    text = extract_text_from_pdf(sample_pdf)
    assert sample_pdf_text in text


def test_extract_text_from_pptx(sample_pptx):
    """Test PPTX text extraction."""
    text = extract_text_from_pptx(sample_pptx)
    assert "Sample PPTX content" in text


def test_process_directory(
    document_processor, learning_resources_dir, sample_pdf, sample_pptx
):
    """Test processing a directory of documents."""
    # Copy sample files to the resources directory
    test_pdf = learning_resources_dir / "test.pdf"
    test_pptx = learning_resources_dir / "test.pptx"
    shutil.copy(sample_pdf, test_pdf)
    shutil.copy(sample_pptx, test_pptx)

    # Mock file processing functions
    with (
        patch("src.document_processor.extract_text_from_pdf") as mock_pdf,
        patch("src.document_processor.extract_text_from_pptx") as mock_pptx,
        patch("openai.OpenAI") as mock_openai_class,
    ):
        mock_pdf.return_value = "PDF content"
        mock_pptx.return_value = "PowerPoint content"
        mock_openai_class.return_value.chat.completions.create.return_value = MagicMock(
            choices=[
                MagicMock(
                    message=MagicMock(
                        content='{"topics": {"Test Topic": {"Subtopic": "Description"}}, "associated_topics": {"Test Topic": ["Related"]}}'
                    )
                )
            ],
            usage=MagicMock(prompt_tokens=100, completion_tokens=50),
        )

        # Set up the OpenAI client for the document processor
        document_processor.openai_helper.client = mock_openai_class.return_value

        # Process the directory
        topics = document_processor.process_directory(learning_resources_dir)

        # Verify the topics structure
        assert isinstance(topics, dict)
        assert "topics" in topics
        assert "Test Topic" in topics["topics"]
        assert mock_pdf.called
        assert mock_pptx.called

        # Verify the mock calls
        mock_pdf.assert_called_with(test_pdf, document_processor.verbosity)
        mock_pptx.assert_called_with(test_pptx, document_processor.verbosity)


def test_process_single_document(document_processor, sample_pdf):
    """Test processing a single document."""
    # Process the document
    topics, cost = document_processor.process_single_document(sample_pdf)

    # Verify the results
    assert isinstance(topics, dict)
    assert "topics" in topics


def test_merge_topics(document_processor):
    """Test merging topics from multiple documents."""
    existing_topics = {
        "topics": {
            "Topic 1": {
                "Subtopic A": {},
                "Subtopic B": {},
            }
        },
        "associated_topics": {"Topic 1": ["Related 1", "Related 2"]},
    }

    new_topics = {
        "topics": {
            "Topic 1": {
                "Subtopic C": {},
            },
            "Topic 2": {
                "Subtopic X": {},
            },
        },
        "associated_topics": {"Topic 1": ["Related 3"], "Topic 2": ["Related X"]},
    }

    # Create a copy of existing topics to merge into
    merged = json.loads(json.dumps(existing_topics))

    # Perform merge
    document_processor._merge_topics(merged, new_topics)

    # Verify merge results
    assert "Topic 1" in merged["topics"]
    assert "Topic 2" in merged["topics"]
    assert "Subtopic A" in merged["topics"]["Topic 1"]
    assert "Subtopic B" in merged["topics"]["Topic 1"]
    assert "Subtopic C" in merged["topics"]["Topic 1"]
    assert set(merged["associated_topics"]["Topic 1"]) == {
        "Related 1",
        "Related 2",
        "Related 3",
    }


def test_existing_topics_awareness(document_processor):
    """Test that existing topics are considered when processing new documents."""
    # Add some existing topics
    document_processor.existing_topics.update(["Topic 1", "Topic 2"])

    # Verify existing topics are tracked
    assert "Topic 1" in document_processor.existing_topics
    assert "Topic 2" in document_processor.existing_topics


def test_openai_helper_token_cost(mock_openai_helper):
    """Test token cost calculation."""
    # Test GPT-4 Turbo costs
    gpt4_cost = mock_openai_helper.calculate_token_cost(
        "gpt-4-turbo-preview", 1000, 500
    )
    assert gpt4_cost == (1000 * 0.01 / 1000) + (500 * 0.03 / 1000)

    # Test unknown model
    unknown_cost = mock_openai_helper.calculate_token_cost("unknown", 1000, 500)
    assert unknown_cost == 0.0


def test_logging_helper(mock_logger):
    """Test logging helper functionality."""
    # Test section header
    mock_logger.log_section_header("Test Section")
    # Test subsection
    mock_logger.log_subsection("Test Subsection")
    # Test token usage
    mock_logger.log_token_usage(100, 50, 0.1, "GPT-4")


@patch("requests.get")
def test_wikipedia_definition(mock_get, note_formatter):
    """Test Wikipedia definition retrieval."""
    # Mock Wikipedia search API response
    mock_search_response = MagicMock()
    mock_search_response.json.return_value = {
        "query": {"search": [{"title": "Test Topic"}]}
    }

    # Mock Wikipedia summary API response
    mock_summary_response = MagicMock()
    mock_summary_response.json.return_value = {"extract": "Test definition"}

    # Configure mock to return different responses for different URLs
    def mock_get_side_effect(url, *args, **kwargs):
        if "action=query" in url:  # Search API
            return mock_search_response
        else:  # Summary API
            return mock_summary_response

    mock_get.side_effect = mock_get_side_effect

    # Get definition
    definition = note_formatter.get_wikipedia_definition("Test Topic")

    # Verify the result
    assert definition == "Test definition"
    assert (
        mock_get.call_count == 2
    )  # Should be called twice - once for search, once for summary


def test_setup_logging(temp_dir):
    """Test logging setup."""
    # Test with different verbosity levels
    for verbosity in [0, 1, 2]:
        setup_logging(verbosity)
        assert logging.root.level == (
            logging.ERROR
            if verbosity == 0
            else logging.INFO
            if verbosity == 1
            else logging.DEBUG
        )

    # Test with log file
    log_file = temp_dir / "test.log"
    setup_logging(1, str(log_file))
    assert log_file.exists()


@patch("src.check_dependencies")
@patch("src.DocumentProcessor")
@patch("builtins.input")
@patch("subprocess.run")
def test_main_interactive_with_edit(
    mock_subprocess_run, mock_input, mock_processor_class, mock_check, temp_dir
):
    """Test the main function in interactive mode with edit functionality."""
    # Create a resource directory for testing
    resource_dir = temp_dir / "Learning Resources"
    resource_dir.mkdir()

    # Setup mock processor
    mock_processor = MagicMock()
    mock_processor.process_directory.return_value = {
        "topics": {"Test Topic": {"Subtopic": {}}},
        "associated_topics": {"Test Topic": ["Related"]},
    }
    mock_processor_class.return_value = mock_processor

    # Mock user input to choose edit option
    mock_input.return_value = "e"

    # Mock nano subprocess to simulate editing
    mock_subprocess_run.return_value = MagicMock(returncode=0)

    # Setup test arguments with interactive mode
    test_args = ["gen_notes.py", str(temp_dir), "-i"]
    with patch("sys.argv", test_args):
        main()

    # Verify dependencies were checked
    mock_check.assert_called_once()

    # Verify directory was processed
    mock_processor.process_directory.assert_called_once_with(Path(resource_dir))

    # Verify notes were created in the parent directory
    notes_dir = temp_dir / "notes"
    assert notes_dir.exists()
    assert (notes_dir / "Test Topic.md").exists()

    # Verify nano was called
    mock_subprocess_run.assert_called_once()
    assert mock_subprocess_run.call_args[0][0][0] == "nano"


def test_document_processor_initialization():
    """Test DocumentProcessor initialization."""
    processor = DocumentProcessor(verbosity=1)
    assert processor.verbosity == 1
    assert isinstance(processor.openai_helper, OpenAIHelper)


def test_process_directory_no_files(tmp_path):
    """Test processing a directory with no files."""
    processor = DocumentProcessor(verbosity=1)
    result = processor.process_directory(tmp_path)
    assert result == {}


def test_process_directory_with_files(
    tmp_path,
    mock_openai_client,
    mock_openai_response,
    sample_pdf,
    sample_pptx,
    sample_mp4,
):
    """Test processing a directory with various file types."""
    # Copy sample files to temp directory
    test_pdf = tmp_path / "test1.pdf"
    test_pptx = tmp_path / "test1.pptx"
    test_mp4 = tmp_path / "test1.mp4"
    shutil.copy(str(sample_pdf), str(test_pdf))
    shutil.copy(str(sample_pptx), str(test_pptx))
    shutil.copy(str(sample_mp4), str(test_mp4))

    # Mock file processing functions
    with (
        patch("src.document_processor.extract_text_from_pdf") as mock_pdf,
        patch("src.document_processor.extract_text_from_pptx") as mock_pptx,
        patch("src.document_processor.transcribe_audio_with_whisper") as mock_mp4,
        patch("openai.OpenAI") as mock_openai_class,
    ):
        mock_pdf.return_value = "PDF content"
        mock_pptx.return_value = "PowerPoint content"
        mock_mp4.return_value = "Video content"
        mock_openai_class.return_value = mock_openai_client

        # Create a processor with the mocked OpenAI client
        processor = DocumentProcessor(verbosity=1)
        processor.openai_helper.client = mock_openai_client

        # Process the directory
        result = processor.process_directory(tmp_path)

        # Verify the results
        assert "topics" in result
        assert "Test Topic" in result["topics"]
        assert mock_pdf.called
        assert mock_pptx.called
        assert mock_mp4.called

        # Verify the mock calls
        mock_pdf.assert_has_calls(
            [
                call(test_pdf, processor.verbosity),
                call(sample_pdf, processor.verbosity),
            ],
            any_order=True,
        )
        mock_pptx.assert_has_calls(
            [
                call(test_pptx, processor.verbosity),
                call(sample_pptx, processor.verbosity),
            ],
            any_order=True,
        )
        mock_mp4.assert_has_calls(
            [
                call(test_mp4, processor.verbosity),
                call(sample_mp4, processor.verbosity),
            ],
            any_order=True,
        )


def test_extract_text_from_pdf(sample_pdf):
    """Test PDF text extraction."""
    with patch("pypdf.PdfReader") as mock_reader:
        mock_page = MagicMock()
        mock_page.extract_text.return_value = "Test content"
        mock_reader.return_value.pages = [mock_page]

        text = extract_text_from_pdf(sample_pdf, verbosity=1)
        assert text == "Test content"


def test_extract_text_from_pptx(sample_pptx):
    """Test PowerPoint text extraction."""
    with patch("pptx.Presentation") as mock_presentation:
        mock_slide = MagicMock()
        mock_shape = MagicMock()
        mock_shape.text = "Test content"
        mock_slide.shapes = [mock_shape]
        mock_presentation.return_value.slides = [mock_slide]

        text = extract_text_from_pptx(sample_pptx, verbosity=1)
        assert text == "Test content"


def test_transcribe_audio_with_whisper(sample_mp4):
    """Test audio transcription."""
    with patch("subprocess.Popen") as mock_popen:
        mock_process = MagicMock()
        mock_process.poll.return_value = 0
        mock_process.returncode = 0
        mock_popen.return_value = mock_process

        # Create a temporary transcript file
        transcript_path = sample_mp4.with_suffix(".txt")
        transcript_path.write_text("Test transcript")

        text = transcribe_audio_with_whisper(sample_mp4, verbosity=1)
        assert text == "Test transcript"
        transcript_path.unlink()


def test_colored_logging():
    """Test colored logging setup."""
    with patch("logging.StreamHandler") as mock_handler:
        setup_colored_logging(1)
        mock_handler.assert_called_once()


def test_progress_bar_verbosity():
    """Test progress bar behavior with different verbosity levels."""
    processor = DocumentProcessor(verbosity=0)
    assert processor.verbosity == 0

    processor = DocumentProcessor(verbosity=1)
    assert processor.verbosity == 1

    processor = DocumentProcessor(verbosity=2)
    assert processor.verbosity == 2


def test_extract_text_from_txt(sample_txt):
    """Test text file extraction."""
    with patch("src.utils.logging") as mock_logging:
        text = extract_text_from_txt(sample_txt)
        assert text == "Sample text content"
        mock_logging.info.assert_called_once()
        mock_logging.error.assert_not_called()


def test_process_directory_with_txt(
    tmp_path,
    mock_openai_client,
    mock_openai_response,
    sample_pdf,
    sample_pptx,
    sample_mp4,
    sample_txt,
):
    """Test processing a directory with various file types including text files."""
    # Copy sample files to temp directory
    test_pdf = tmp_path / "test1.pdf"
    test_pptx = tmp_path / "test1.pptx"
    test_mp4 = tmp_path / "test1.mp4"
    test_txt = tmp_path / "test1.txt"
    shutil.copy(str(sample_pdf), str(test_pdf))
    shutil.copy(str(sample_pptx), str(test_pptx))
    shutil.copy(str(sample_mp4), str(test_mp4))
    shutil.copy(str(sample_txt), str(test_txt))

    # Mock file processing functions
    with (
        patch("src.document_processor.extract_text_from_pdf") as mock_pdf,
        patch("src.document_processor.extract_text_from_pptx") as mock_pptx,
        patch("src.document_processor.transcribe_audio_with_whisper") as mock_mp4,
        patch("src.document_processor.extract_text_from_txt") as mock_txt,
        patch("openai.OpenAI") as mock_openai_class,
        patch("src.document_processor.logging") as mock_logging,
        patch("src.openai_helper.logging") as mock_openai_logging,
        patch("src.utils.logging") as mock_utils_logging,
    ):
        mock_pdf.return_value = "PDF content"
        mock_pptx.return_value = "PowerPoint content"
        mock_mp4.return_value = "Video content"
        mock_txt.return_value = "Text content"
        mock_openai_class.return_value = mock_openai_client

        # Create a processor with the mocked OpenAI client
        processor = DocumentProcessor(verbosity=1)
        processor.openai_helper.client = mock_openai_client

        # Process the directory
        result = processor.process_directory(tmp_path)

        # Verify the results
        assert "topics" in result
        assert "Test Topic" in result["topics"]
        assert mock_pdf.called
        assert mock_pptx.called
        assert mock_mp4.called
        assert mock_txt.called

        # Verify the mock calls
        mock_pdf.assert_has_calls(
            [
                call(test_pdf, processor.verbosity),
                call(sample_pdf, processor.verbosity),
            ],
            any_order=True,
        )
        mock_pptx.assert_has_calls(
            [
                call(test_pptx, processor.verbosity),
                call(sample_pptx, processor.verbosity),
            ],
            any_order=True,
        )
        mock_mp4.assert_has_calls(
            [
                call(test_mp4, processor.verbosity),
                call(sample_mp4, processor.verbosity),
            ],
            any_order=True,
        )
        mock_txt.assert_has_calls(
            [
                call(test_txt, processor.verbosity),
                call(sample_txt, processor.verbosity),
            ],
            any_order=True,
        )


def test_clean_json_response():
    """Test cleaning of improperly formatted JSON responses."""
    from src.openai_helper import OpenAIHelper

    helper = OpenAIHelper()

    # Test case 1: JSON wrapped in markdown code block
    input_json = """```json
{
    "topics": {
        "Test Topic": {
            "Subtopic": "Content"
        }
    }
}
```"""
    expected = '{\n    "topics": {\n        "Test Topic": {\n            "Subtopic": "Content"\n        }\n    }\n}'
    assert helper._clean_json_response(input_json) == expected

    # Test case 2: JSON with extra text before and after
    input_json = 'Some text before {"topics": {"Test": "Content"}} some text after'
    expected = '{"topics": {"Test": "Content"}}'
    assert helper._clean_json_response(input_json) == expected

    # Test case 3: JSON with multiple code blocks
    input_json = """First block
```json
{
    "topics": {
        "Test": "Content"
    }
}
```
Second block"""
    expected = '{\n    "topics": {\n        "Test": "Content"\n    }\n}'
    assert helper._clean_json_response(input_json) == expected

    # Test case 4: Invalid JSON (no braces)
    with pytest.raises(ValueError, match="No valid JSON object found in response"):
        helper._clean_json_response("This is not a JSON object")


def test_generate_topics_with_improper_json():
    """Test topic generation with improperly formatted JSON responses."""
    from src.openai_helper import OpenAIHelper
    from unittest.mock import MagicMock, patch

    # Mock logging to prevent TypeError
    with patch("logging.info"), patch("logging.error"), patch("logging.debug"):
        helper = OpenAIHelper()
        helper.add_document("Test content")

        # Mock response with markdown-wrapped JSON
        mock_response = MagicMock()
        mock_response.choices = [
            MagicMock(
                message=MagicMock(
                    content="""```json
{
    "topics": {
        "Test Topic": {
            "Subtopic": "Content"
        }
    },
    "associated_topics": {
        "Test Topic": ["Related Topic"]
    }
}
```"""
                )
            )
        ]
        mock_response.usage = MagicMock(prompt_tokens=100, completion_tokens=50)

        with patch.object(
            helper.client.chat.completions, "create", return_value=mock_response
        ):
            topics = helper.generate_topics(verbosity=2)

            assert "Test Topic" in topics["topics"]
            assert "Subtopic" in topics["topics"]["Test Topic"]
            assert topics["topics"]["Test Topic"]["Subtopic"] == "Content"
            assert "Test Topic" in topics["associated_topics"]
            assert "Related Topic" in topics["associated_topics"]["Test Topic"]


def test_document_chunking():
    """Test document chunking functionality."""
    from src.openai_helper import OpenAIHelper

    helper = OpenAIHelper()

    # Test case 1: Small document (should not be chunked)
    small_doc = "This is a small document.\n\nIt has two paragraphs."
    chunks = helper._chunk_document(small_doc)
    assert len(chunks) == 1
    assert chunks[0] == small_doc

    # Set a small chunk size for testing
    helper.chunk_size = 100  # much smaller than default

    # Test case 2: Large document (should be chunked)
    large_doc = "\n\n".join(["Paragraph " + str(i) for i in range(1000)])
    chunks = helper._chunk_document(large_doc)
    assert len(chunks) > 1

    # Test case 3: Document with very large paragraphs
    large_paragraph = "x" * (
        helper.chunk_size * 4
    )  # Create a paragraph larger than chunk size
    doc_with_large_para = (
        f"Small paragraph.\n\n{large_paragraph}\n\nAnother small paragraph."
    )
    chunks = helper._chunk_document(doc_with_large_para)
    assert len(chunks) > 1
    assert all(len(chunk) <= helper.chunk_size * 4 for chunk in chunks)


def test_token_rate_limiting():
    """Test token rate limiting functionality."""
    from src.openai_helper import OpenAIHelper
    import time

    helper = OpenAIHelper()

    # Test case 1: Token usage tracking
    helper.token_usage_history = [(time.time(), 1000)]
    current_usage = helper._get_current_token_usage()
    assert current_usage == 1000

    # Test case 2: Token history cleaning
    old_time = time.time() - 61  # Just over a minute old
    helper.token_usage_history = [(old_time, 1000), (time.time(), 500)]
    helper._clean_token_history()
    assert len(helper.token_usage_history) == 1
    assert helper.token_usage_history[0][1] == 500


def test_api_call_with_retry():
    """Test API call retry logic."""
    from src.openai_helper import OpenAIHelper
    from unittest.mock import patch, MagicMock

    helper = OpenAIHelper()

    # Mock the OpenAI client
    mock_response = MagicMock()
    mock_response.choices = [
        MagicMock(message=MagicMock(content='{"topics": {"Test": "Content"}}'))
    ]
    mock_response.usage = MagicMock(prompt_tokens=100, completion_tokens=50)

    with patch.object(helper.client.chat.completions, "create") as mock_create:
        # Test case 1: Successful call
        mock_create.return_value = mock_response
        response = helper._make_api_call_with_retry(
            [{"role": "user", "content": "test"}]
        )
        assert response == mock_response
        assert mock_create.call_count == 1

        # Test case 2: Rate limit with retry
        mock_create.side_effect = [Exception("rate_limit_exceeded"), mock_response]
        response = helper._make_api_call_with_retry(
            [{"role": "user", "content": "test"}]
        )
        assert response == mock_response
        assert mock_create.call_count == 3  # Initial + 2 retries

        # Test case 3: Non-rate-limit error
        mock_create.side_effect = Exception("other_error")
        with pytest.raises(Exception, match="other_error"):
            helper._make_api_call_with_retry([{"role": "user", "content": "test"}])


if __name__ == "__main__":
    pytest.main(["-v"])
