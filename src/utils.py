"""Utility functions for the AI Note Workflow package."""

import sys
import subprocess
import logging
from pathlib import Path
from typing import Optional
from tqdm import tqdm
import time
import os
from .logging_helper import (
    setup_colored_logging,
    create_progress_bar,
    update_progress_bar,
    close_progress_bar,
)


def setup_logging(verbosity: int, log_file: str = None):
    """Set up logging based on verbosity level and optional log file."""
    setup_colored_logging(verbosity, log_file)


def check_dependencies():
    """Check if required dependencies are installed."""
    try:
        import pypdf
        import pptx
    except ImportError as e:
        logging.error(f"Required dependency not installed: {e}")
        logging.info("Please install required packages: pip install pypdf python-pptx")
        sys.exit(1)


def extract_text_from_pdf(pdf_path: Path, verbosity: int = 1) -> str:
    """Extract text from a PDF file."""
    try:
        from pypdf import PdfReader

        logging.info(f"Processing PDF: {pdf_path.name}")
        reader = PdfReader(str(pdf_path))
        total_pages = len(reader.pages)

        # Create progress bar
        pbar = create_progress_bar(total_pages, "Extracting PDF text", verbosity)

        text_runs = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text_runs.append(text)
            update_progress_bar(pbar, i + 1, total_pages)

        close_progress_bar(pbar)
        return "\n".join(text_runs)
    except Exception as e:
        logging.error(f"Error processing PDF {pdf_path.name}: {e}")
        return ""


def extract_text_from_pptx(pptx_path: Path, verbosity: int = 1) -> str:
    """Extract text from a PowerPoint file."""
    try:
        from pptx import Presentation

        logging.info(f"Processing PowerPoint: {pptx_path.name}")
        prs = Presentation(str(pptx_path))
        total_slides = len(prs.slides)

        # Create progress bar
        pbar = create_progress_bar(
            total_slides, "Extracting PowerPoint text", verbosity
        )

        text_runs = []
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
            update_progress_bar(pbar, i + 1, total_slides)

        close_progress_bar(pbar)
        return "\n".join(text_runs)
    except Exception as e:
        logging.error(f"Error processing PPTX {pptx_path.name}: {e}")
        return ""


def transcribe_audio_with_whisper(mp4_path: Path, verbosity: int = 1) -> str:
    """Transcribe audio from an MP4 file using Whisper."""
    try:
        logging.info(f"Transcribing {mp4_path.name} using Whisper...")

        # Get file size for progress estimation
        file_size = os.path.getsize(mp4_path)

        # Create progress bar
        pbar = create_progress_bar(100, "Transcribing audio", verbosity)

        # Start the Whisper process
        process = subprocess.Popen(
            [
                "whisper",
                str(mp4_path),
                "--model",
                "base",
                "--language",
                "en",
                "--output_format",
                "txt",
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )

        # Update progress bar while process is running
        while process.poll() is None:
            # Estimate progress based on output file size
            if os.path.exists(mp4_path.with_suffix(".txt")):
                current_size = os.path.getsize(mp4_path.with_suffix(".txt"))
                update_progress_bar(pbar, current_size, file_size)
            time.sleep(0.1)

        close_progress_bar(pbar)

        # Check for errors
        if process.returncode != 0:
            error = process.stderr.read()
            logging.error(f"Error running Whisper on {mp4_path.name}: {error}")
            return ""

        transcript_file = mp4_path.with_suffix(".txt")
        return transcript_file.read_text() if transcript_file.exists() else ""
    except Exception as e:
        logging.error(f"Error processing MP4 {mp4_path.name}: {e}")
        return ""
